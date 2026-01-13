#!/usr/bin/env python3
"""
Generate PowerPoint presentation from the hackathon script.
Run: python3 create_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme
PRIMARY_COLOR = RGBColor(59, 130, 246)  # Blue
ACCENT_COLOR = RGBColor(34, 197, 94)   # Green
TEXT_COLOR = RGBColor(30, 30, 30)      # Dark gray
LIGHT_BG = RGBColor(245, 245, 245)     # Light gray

def add_title_slide(title, subtitle=""):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_COLOR
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1.5))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True
        p = subtitle_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(28)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

def add_content_slide(title, content_list, accent_title=True):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR if accent_title else TEXT_COLOR
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_list):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(24)
        p.font.color.rgb = TEXT_COLOR
        p.level = 0
        p.space_before = Pt(12)

def add_two_column_slide(title, left_content, right_content):
    """Add a two-column content slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5.5))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    for i, item in enumerate(left_content):
        if i > 0:
            left_frame.add_paragraph()
        p = left_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(8)
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.3), Inches(5.5))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    for i, item in enumerate(right_content):
        if i > 0:
            right_frame.add_paragraph()
        p = right_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(8)

# SLIDE 1: Title
add_title_slide("🤖 AI PR Code Reviewer", "Making Code Review Instant & Intelligent")

# SLIDE 2: The Problem
add_content_slide("The Problem", [
    "⏳ Code reviews are slow - senior engineers are bottlenecked",
    "😕 They're inconsistent - different reviewers catch different things",
    "😰 They're overwhelming - 50-file PRs? Good luck catching bugs",
    "📉 Result: slower development, missed security issues"
])

# SLIDE 3: The Solution
add_content_slide("Our Solution", [
    "🎯 Instant AI-powered code analysis - seconds, not days",
    "🔍 Identifies security, performance & maintainability issues",
    "📋 Auto-generates comprehensive test plans",
    "📊 Merge readiness score (0-100)",
    "💬 Optional GitHub integration - post reviews as comments"
])

# SLIDE 4: How It Works
add_two_column_slide("How It Works", 
    [
        "1. User submits PR URL",
        "2. Backend fetches from GitHub",
        "   - PR metadata",
        "   - Full diff",
        "   - Changed files",
        "",
        "3. AI Analyzer processes",
    ],
    [
        "4. Structured results:",
        "   - Summary",
        "   - Findings",
        "   - Risk matrix",
        "   - Test plan",
        "",
        "5. Display in UI / Post to GitHub"
    ]
)

# SLIDE 5: Key Features
add_content_slide("Key Features", [
    "✨ Dual-mode analysis: AI (with OpenAI) + Heuristics (always works)",
    "🔐 Security scanning: finds secrets, SQL injection, XSS",
    "⚡ Performance analysis: detects N+1 queries, nested loops",
    "🧪 Test suggestions: unit, integration, edge cases",
    "📈 Confidence scores: know how certain we are about each finding"
])

# SLIDE 6: Tech Stack
add_two_column_slide("Tech Stack",
    [
        "Backend:",
        "• FastAPI (Python)",
        "• SQLAlchemy ORM",
        "• Async HTTP (httpx)",
        "• GitHub API integration",
    ],
    [
        "Frontend:",
        "• React + TypeScript",
        "• Vite",
        "• Tailwind CSS",
        "• shadcn/ui components",
    ]
)

# SLIDE 7: AI Layer
add_content_slide("The AI Magic", [
    "🧠 OpenAI GPT-4o-mini with strict JSON schema",
    "✅ Structured output validation with Pydantic",
    "🔄 Retry logic for malformed responses",
    "⚙️ Fallback to heuristics if AI fails",
    "⚡ ~5-10 seconds per analysis"
])

# SLIDE 8: Demo Setup
add_content_slide("Live Demo", [
    "Let's analyze a real Pull Request...",
    "",
    "📍 Frontend: https://colorstackwinterhackathon2025-codes.vercel.app",
    "📍 Backend: https://colorstackwinterhackathon2025-codesearch.onrender.com",
    "",
    "[Demo walkthrough on next slides]"
])

# SLIDE 9: Demo - Submit
add_content_slide("Demo Step 1: Submit PR", [
    "✓ Enter GitHub PR URL",
    "✓ Click submit",
    "✓ Backend fetches PR from GitHub",
    "✓ Analyzer processes changes",
])

# SLIDE 10: Demo - Results
add_content_slide("Demo Step 2: View Results", [
    "✓ PR Summary - what changed & why",
    "✓ Prioritized Findings - ranked by severity",
    "✓ Risk Matrix - 4D assessment",
    "✓ Test Plan - actionable suggestions",
    "✓ Merge Readiness Score - 0-100"
])

# SLIDE 11: Challenges
add_content_slide("Challenges We Solved", [
    "🔴 CORS preflight hell - 400 errors between Vercel & Render",
    "→ Solution: Dynamic ALLOWED_ORIGINS with regex patterns",
    "",
    "🔴 LLM output validation - JSON malformed sometimes",
    "→ Solution: Schema validation + retry + fallback",
    "",
    "🔴 Large diffs exceeding token limits",
    "→ Solution: Truncate intelligently, prioritize changed lines"
])

# SLIDE 12: Technical Highlights
add_content_slide("Technical Highlights", [
    "🚀 Concurrent GitHub API calls with asyncio",
    "🎯 Weighted scoring algorithm (critical = 30pts, high = 15pts)",
    "🔄 Dual-mode fallback system",
    "📊 JSONB storage for analysis results",
    "🌍 Production deployment on Render + Vercel"
])

# SLIDE 13: What We Built
add_content_slide("What We Accomplished", [
    "✅ 12,000+ lines of production-ready code",
    "✅ Full-stack web app (backend + frontend)",
    "✅ Real GitHub integration (fetching & posting)",
    "✅ AI + heuristics analysis engine",
    "✅ Deployed to production (live now!)"
])

# SLIDE 14: What's Next
add_two_column_slide("What's Next",
    [
        "Near-term:",
        "• CI/CD integration",
        "• Custom rules engine",
        "• Line-specific comments",
        "• Historical tracking"
    ],
    [
        "Long-term:",
        "• GitHub App webhooks",
        "• Multi-language linters",
        "• Team dashboards",
        "• Self-hosted option"
    ]
)

# SLIDE 15: Closing
add_title_slide("Code review shouldn't be a bottleneck.\nIt should be instant, comprehensive & accessible.", "AI PR Code Reviewer")

# SLIDE 16: Questions
add_title_slide("❓ Questions?", "github.com/Johnnti/colorstackwinterhackathon2025-codesearch")

# Save presentation
output_file = "AI_PR_Code_Reviewer_Presentation.pptx"
prs.save(output_file)
print(f"✅ Presentation created: {output_file}")
print(f"📍 Location: /Users/john/colorstackwinterhackathon2025-codesearch/{output_file}")
